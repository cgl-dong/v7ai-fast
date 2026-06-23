"""Skill management API — list and preview document transformation skills."""
import io
from typing import Optional, List
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Query, BackgroundTasks
from pydantic import BaseModel
from sqlalchemy.orm import Session

from app.core.database import get_db, User, KnowledgeFile, ChatSession, ChatMessage
from app.api.v1.endpoints.auth import get_current_user
from app.api.v1.endpoints.knowledge import get_optional_user
from app.services.skill_base import registry
from app.services.knowledge import KnowledgeService
from app.services.indexer import Indexer
from app.core.logging import logger

# Trigger skill auto-discovery on module import
import app.services.skills  # noqa: F401 — ensures skills are registered before first API call

router = APIRouter()


class SkillInfo(BaseModel):
    name: str
    description: str
    input_types: List[str]
    output_type: str


class PreviewRequest(BaseModel):
    file_id: int
    skill_names: List[str] = []  # ordered list of skills to apply


class SkillApplyRequest(BaseModel):
    """Apply skills to a file and re-index with the transformed content."""
    skill_names: List[str]  # ordered list of skills to apply


# ── Skill listing ──────────────────────────────────────────────────


@router.get("")
async def list_skills():
    """List all available skills (transform + tool)."""
    skills = registry.list_skills()
    return {
        "skills": [s.to_dict() for s in skills],
        "count": len(skills),
    }


# ── Skill preview ─────────────────────────────────────────────────


@router.post("/preview")
async def preview_skill(
    req: PreviewRequest,
    db: Session = Depends(get_db),
    current_user: Optional[User] = Depends(get_optional_user),
):
    """Preview what a skill pipeline would produce for a file.

    Applies the skill chain to the stored file content, parses the result,
    and returns the first N characters of the transformed text.
    """
    svc = KnowledgeService(db)
    user_id = current_user.id if current_user else None
    record = svc.get_file_by_id(req.file_id, user_id=user_id)
    if not record:
        raise HTTPException(status_code=404, detail="文件不存在")

    try:
        content_bytes, _, _ = svc.get_file_content(req.file_id)
    except FileNotFoundError as e:
        raise HTTPException(status_code=404, detail=str(e))

    # Run skill pipeline
    try:
        new_content, new_filename, new_file_type = registry.run_pipeline(
            content=content_bytes,
            filename=record.filename,
            file_type=record.file_type,
            skill_names=req.skill_names,
        )
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except RuntimeError as e:
        raise HTTPException(status_code=500, detail=str(e))

    # Parse the transformed content
    idx = Indexer(db)
    try:
        text = idx._parse_content(new_content, new_filename, new_file_type)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"解析转换后文件失败: {e}")

    return {
        "file_id": req.file_id,
        "original_filename": record.filename,
        "original_type": record.file_type,
        "original_size": record.file_size,
        "skills_applied": req.skill_names,
        "transformed_filename": new_filename,
        "transformed_type": new_file_type,
        "transformed_size": len(new_content),
        "preview_text": text[:5000],
        "preview_length": min(len(text), 5000),
        "total_length": len(text),
    }


# ── Apply skills and re-index ─────────────────────────────────────


@router.post("/files/{file_id}/apply-skills")
async def apply_skills_to_file(
    file_id: int,
    req: SkillApplyRequest,
    db: Session = Depends(get_db),
    background_tasks: BackgroundTasks = None,
    current_user: User = Depends(get_current_user),
):
    """Apply a skill pipeline to a file and trigger re-indexing.

    The original file in MinIO is unchanged. The skill pipeline runs
    in-memory during indexing to produce the final text for embedding.

    Example: POST /api/v1/files/42/apply-skills
             {"skill_names": ["pdf_to_docx"]}
    """
    from app.api.v1.endpoints.knowledge import _run_index_task

    record = db.query(KnowledgeFile).filter(KnowledgeFile.id == file_id).first()
    if not record:
        raise HTTPException(status_code=404, detail="文件不存在")

    if record.user_id is not None and record.user_id != current_user.id:
        raise HTTPException(status_code=403, detail="无权限操作此文件")

    # Validate skills exist
    for name in req.skill_names:
        if registry.get(name) is None:
            available = [s.name for s in registry.list_skills()]
            raise HTTPException(
                status_code=400,
                detail=f"技能 '{name}' 不存在。可用技能: {available}",
            )

    # Validate pipeline compatibility
    try:
        registry.build_pipeline(record.file_type, req.skill_names)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))

    # Launch background index with skills
    background_tasks.add_task(
        _run_index_task, file_id,
        strategy=getattr(record, 'chunk_strategy', None),
        chunk_size=getattr(record, 'chunk_size', None),
        chunk_overlap=getattr(record, 'chunk_overlap', None),
        skills=req.skill_names,
    )

    logger.info(f"Skill index task queued: file={file_id}, skills={req.skill_names}")
    return {
        "message": "技能管线已提交，正在后台处理",
        "file_id": file_id,
        "skills": req.skill_names,
    }


# ── Chat skill execution ──────────────────────────────────────────


class ChatSkillExecuteRequest(BaseModel):
    """Execute a tool skill on chat conversation content."""
    skill_name: str  # e.g., "ppt-generation"
    session_id: str  # chat session ID


@router.post("/execute")
async def execute_skill_on_chat(
    req: ChatSkillExecuteRequest,
    db: Session = Depends(get_db),
    current_user: Optional[User] = Depends(get_optional_user),
):
    """Execute a tool skill on a chat session's conversation content.

    Fetches all messages from the session, then invokes the skill
    with the conversation as input. Returns a downloadable file.

    Supported skills:
    - ppt-generation: exports conversation as a PPTX presentation
    """
    # Validate skill exists
    info = registry.get_info(req.skill_name)
    if info is None:
        available = [s.name for s in registry.list_skills()]
        raise HTTPException(
            status_code=400,
            detail=f"技能 '{req.skill_name}' 不存在。可用: {available}",
        )

    # Fetch session messages
    session = db.query(ChatSession).filter(
        ChatSession.chat_id == req.session_id
    ).first()
    if not session:
        raise HTTPException(status_code=404, detail="会话不存在")

    messages = db.query(ChatMessage).filter(
        ChatMessage.session_id == session.id
    ).order_by(ChatMessage.created_at).all()

    if not messages:
        raise HTTPException(status_code=400, detail="会话无消息内容")

    # Build conversation text
    conv_lines = []
    for msg in messages:
        role_label = "用户" if msg.role == "user" else "AI助手"
        conv_lines.append(f"【{role_label}】{msg.content or ''}")
    conversation = "\n\n".join(conv_lines)

    # Execute skill based on type
    if req.skill_name == "ppt-generation":
        return _execute_ppt_generation(conversation, session.title or session.chat_id)
    else:
        # Generic tool skill — return conversation content for external processing
        return {
            "skill": req.skill_name,
            "session_id": req.session_id,
            "message_count": len(messages),
            "content": conversation,
            "message": f"技能 '{req.skill_name}' 已准备好输入数据",
        }


def _execute_ppt_generation(conversation: str, title: str) -> dict:
    """Generate a simple PPTX from conversation text using python-pptx."""
    import tempfile
    import os
    import uuid
    import base64

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="python-pptx 未安装。请运行: pip install python-pptx",
        )

    # Split conversation into slide-sized chunks
    paragraphs = [p.strip() for p in conversation.split("\n\n") if p.strip()]
    slides_data = []
    current_slide = ""
    for p in paragraphs:
        if len(current_slide) + len(p) < 800:
            current_slide += p + "\n\n"
        else:
            if current_slide:
                slides_data.append(current_slide.strip())
            current_slide = p + "\n\n"
    if current_slide.strip():
        slides_data.append(current_slide.strip())

    if not slides_data:
        slides_data = [conversation[:2000]]

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    txBox = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title[:80] if title else "对话导出"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x66, 0x7E, 0xEA)
    p.alignment = 1  # center

    subtitle = title_slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
    stf = subtitle.text_frame
    sp = stf.paragraphs[0]
    sp.text = f"共 {len(slides_data)} 页 · 由 v7ai-fast 技能生成"
    sp.font.size = Pt(18)
    sp.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    sp.alignment = 1

    # Content slides
    colors = [
        RGBColor(0x66, 0x7E, 0xEA),
        RGBColor(0x52, 0xC4, 0x1A),
        RGBColor(0xF6, 0xAD, 0x55),
        RGBColor(0xFF, 0x4D, 0x4F),
        RGBColor(0x72, 0x2E, 0xD1),
    ]

    for i, text in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Header bar
        left, top, width, height = Inches(0), Inches(0), Inches(13.333), Inches(1.2)
        shape = slide.shapes.add_shape(1, left, top, width, height)  # rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i % len(colors)]
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"第 {i + 1} 页"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True

        # Content
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True

        lines = text.split("\n")
        for j, line in enumerate(lines):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line[:200]
            if line.startswith("【"):
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = colors[i % len(colors)]
            else:
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.space_after = Pt(8)

        # Page number
        footer = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4))
        ftf = footer.text_frame
        fp = ftf.paragraphs[0]
        fp.text = f"{i + 1}/{len(slides_data)}"
        fp.font.size = Pt(10)
        fp.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        fp.alignment = 2  # right

    # Save to temp file and encode
    fd, tmp_path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    try:
        prs.save(tmp_path)
        with open(tmp_path, "rb") as f:
            pptx_bytes = f.read()
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass

    # Save to MinIO for download
    from minio import Minio
    from app.core.settings import settings

    safe_title = title[:30].replace("/", "_").replace("\\", "_") if title else "conversation"
    stored_name = f"exports/{uuid.uuid4().hex}.pptx"
    client = Minio(
        endpoint=settings.minio_endpoint,
        access_key=settings.minio_access_key,
        secret_key=settings.minio_secret_key,
        secure=settings.minio_secure,
    )
    if not client.bucket_exists(settings.minio_bucket):
        client.make_bucket(settings.minio_bucket)
    client.put_object(
        bucket_name=settings.minio_bucket,
        object_name=stored_name,
        data=io.BytesIO(pptx_bytes),
        length=len(pptx_bytes),
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    # Generate presigned download URL (24h expiry)
    download_url = client.presigned_get_object(
        settings.minio_bucket, stored_name, expires=86400
    )

    logger.info(
        f"PPT generated: {len(slides_data)} slides, "
        f"size={len(pptx_bytes)}B, session={title}"
    )
    return {
        "skill": "ppt-generation",
        "slides": len(slides_data),
        "size_bytes": len(pptx_bytes),
        "download_url": download_url,
        "filename": f"{safe_title}.pptx",
        "message": f"已生成 {len(slides_data)} 页演示文稿",
    }
